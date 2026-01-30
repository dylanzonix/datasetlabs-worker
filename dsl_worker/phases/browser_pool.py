"""
Browser Pool for concurrent browser automation.

Manages a pool of browser instances for parallel tasks.
Supports:
- Simple fetch (no LLM, just grab page content)
- Browse with checkpoints (research agent stays in control)
"""

import asyncio
import logging
import os
import uuid
import hashlib
from contextlib import asynccontextmanager
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Awaitable

from markdownify import markdownify as md

logger = logging.getLogger(__name__)


@dataclass
class BrowseResult:
    """Result from a browser task."""
    success: bool
    pages_marked: int = 0
    final_url: Optional[str] = None
    error: Optional[str] = None


def url_to_filename(url: str) -> str:
    """Convert URL to safe filename."""
    url_hash = hashlib.md5(url.encode()).hexdigest()[:8]
    from urllib.parse import urlparse
    parsed = urlparse(url)
    domain = parsed.netloc.replace(".", "_").replace(":", "_")
    return f"{domain}_{url_hash}.md"


class BrowserPool:
    """
    Pool of browser instances for parallel automation.
    """
    
    def __init__(
        self,
        size: int = 5,
        profiles_dir: str = "./browser_profiles",
        headless: bool = False,
    ):
        self.size = size
        self.profiles_dir = Path(profiles_dir)
        self.profiles_dir.mkdir(parents=True, exist_ok=True)
        self.headless = headless
        
        self._browsers: List[Any] = []
        self._browser_ids: Dict[int, int] = {}
        self._available: asyncio.Queue = asyncio.Queue()
        self._semaphore = asyncio.Semaphore(size)
        self._started = False
        self._start_lock = asyncio.Lock()
        
        logger.info(f"[BrowserPool] Initialized: size={size}, headless={headless}")
    
    async def start(self):
        """Initialize all browser instances."""
        async with self._start_lock:
            if self._started:
                return
            
            from browser_use import Browser
            
            logger.info(f"[BrowserPool] Starting {self.size} browsers...")
            
            for i in range(self.size):
                profile_dir = self.profiles_dir / f"browser-{i}"
                profile_dir.mkdir(exist_ok=True)
                
                browser = Browser(
                    user_data_dir=str(profile_dir),
                    headless=self.headless,
                )
                await browser.start()
                
                self._browsers.append(browser)
                self._browser_ids[id(browser)] = i
                await self._available.put(browser)
                
                logger.info(f"[BrowserPool] Browser {i} started")
            
            self._started = True
            logger.info(f"[BrowserPool] All {self.size} browsers ready")
    
    async def stop(self):
        """Shutdown all browsers."""
        for browser in self._browsers:
            try:
                await browser.stop()
            except Exception as e:
                logger.warning(f"[BrowserPool] Error stopping browser: {e}")
        
        self._browsers = []
        self._available = asyncio.Queue()
        self._started = False
        logger.info("[BrowserPool] Stopped")
    
    @asynccontextmanager
    async def acquire(self):
        """Acquire a browser from the pool."""
        async with self._semaphore:
            browser = await self._available.get()
            browser_id = self._browser_ids.get(id(browser), "?")
            logger.debug(f"[BrowserPool] Acquired browser {browser_id}")
            try:
                yield browser
            finally:
                await self._available.put(browser)
                logger.debug(f"[BrowserPool] Released browser {browser_id}")
    
    # =========================================================================
    # Simple Fetch (no LLM)
    # =========================================================================
    
    DOWNLOADABLE_EXTENSIONS = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'}
    ARCHIVE_EXTENSIONS = {'.zip', '.rar', '.gz', '.tar'}
    
    async def simple_fetch(self, url: str, wait_time: float = 1.5) -> str:
        """
        Fast fetch - no LLM, just browser + JS rendering.
        Returns markdown for HTML pages, or downloads and parses files.
        """
        url_lower = url.lower()
        
        # Skip archives
        for ext in self.ARCHIVE_EXTENSIONS:
            if ext in url_lower:
                logger.warning(f"[BrowserPool] Skipping archive: {url[:60]}...")
                return f"[Cannot process archive file: {url}]"
        
        # Handle downloadable files
        for ext in self.DOWNLOADABLE_EXTENSIONS:
            if ext in url_lower:
                logger.info(f"[BrowserPool] Downloading file: {url[:60]}...")
                return await self._download_and_parse_file(url, ext)
        
        # Regular HTML page
        async with self.acquire() as browser:
            browser_id = self._browser_ids.get(id(browser), "?")
            page = None
            try:
                logger.info(f"[Browser-{browser_id}] Fetching: {url[:80]}...")
                page = await browser.new_page(url)
                await asyncio.sleep(wait_time)
                
                html = await page.evaluate('() => document.body.innerHTML')
                markdown = md(html, heading_style='ATX')
                
                logger.info(f"[Browser-{browser_id}] Got {len(markdown)} chars")
                return markdown
                
            except Exception as e:
                logger.error(f"[Browser-{browser_id}] Fetch failed: {e}")
                return f"Error fetching {url}: {e}"
            finally:
                if page:
                    try:
                        await browser.close_page(page)
                    except:
                        pass
    
    async def _download_and_parse_file(self, url: str, ext: str) -> str:
        """Download a file and extract text content."""
        import tempfile
        import httpx
        
        try:
            async with httpx.AsyncClient(follow_redirects=True, timeout=60.0) as client:
                response = await client.get(url)
                response.raise_for_status()
                content = response.content
            
            logger.info(f"[BrowserPool] Downloaded {len(content)} bytes")
            
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as f:
                f.write(content)
                temp_path = f.name
            
            try:
                if ext == '.pdf':
                    return await self._parse_pdf(temp_path, url)
                elif ext in ('.docx', '.doc'):
                    return await self._parse_docx(temp_path, url)
                elif ext in ('.xlsx', '.xls'):
                    return await self._parse_excel(temp_path, url)
                elif ext in ('.pptx', '.ppt'):
                    return await self._parse_pptx(temp_path, url)
                else:
                    return f"[Downloaded but cannot parse {ext}: {url}]"
            finally:
                import os
                try:
                    os.unlink(temp_path)
                except:
                    pass
                    
        except Exception as e:
            logger.error(f"[BrowserPool] Download failed: {e}")
            return f"[Error downloading {url}: {e}]"
    
    async def _parse_pdf(self, path: str, url: str) -> str:
        """Extract text from PDF."""
        try:
            import pdfplumber
            
            text_parts = []
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages[:50]):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {i+1} ---\n{page_text}")
            
            if text_parts:
                return f"# PDF: {url}\n\n" + "\n\n".join(text_parts)
            return f"[PDF at {url} contains no extractable text]"
            
        except ImportError:
            try:
                from pypdf import PdfReader
                reader = PdfReader(path)
                text_parts = []
                for i, page in enumerate(reader.pages[:50]):
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"--- Page {i+1} ---\n{text}")
                if text_parts:
                    return f"# PDF: {url}\n\n" + "\n\n".join(text_parts)
                return f"[PDF at {url} contains no extractable text]"
            except ImportError:
                return f"[Cannot parse PDF - install pdfplumber: {url}]"
        except Exception as e:
            return f"[Error parsing PDF {url}: {e}]"
    
    async def _parse_docx(self, path: str, url: str) -> str:
        """Extract text from DOCX."""
        try:
            from docx import Document
            doc = Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            if paragraphs:
                return f"# DOCX: {url}\n\n" + "\n\n".join(paragraphs)
            return f"[DOCX at {url} contains no text]"
        except ImportError:
            return f"[Cannot parse DOCX - install python-docx: {url}]"
        except Exception as e:
            return f"[Error parsing DOCX {url}: {e}]"
    
    async def _parse_excel(self, path: str, url: str) -> str:
        """Extract text from Excel."""
        try:
            import pandas as pd
            sheets = pd.read_excel(path, sheet_name=None)
            parts = [f"# Excel: {url}\n"]
            for sheet_name, df in sheets.items():
                parts.append(f"\n## Sheet: {sheet_name}\n")
                parts.append(df.to_markdown(index=False))
            return "\n".join(parts)
        except ImportError:
            return f"[Cannot parse Excel - install pandas openpyxl: {url}]"
        except Exception as e:
            return f"[Error parsing Excel {url}: {e}]"
    
    async def _parse_pptx(self, path: str, url: str) -> str:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_parts = [f"--- Slide {i+1} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text)
                if len(slide_parts) > 1:
                    slides_text.append("\n".join(slide_parts))
            if slides_text:
                return f"# PowerPoint: {url}\n\n" + "\n\n".join(slides_text)
            return f"[PPTX at {url} contains no text]"
        except ImportError:
            return f"[Cannot parse PPTX - install python-pptx: {url}]"
        except Exception as e:
            return f"[Error parsing PPTX {url}: {e}]"
    
    # =========================================================================
    # Browse with Checkpoints (Research agent stays in control)
    # =========================================================================
    
    async def browse_with_checkpoints(
        self,
        url: str,
        goal: str,
        llm,
        checkpoint_callback: Callable[[str, str], Awaitable[str]],
        extraction_dir: Path,
        extraction_queue: asyncio.Queue,
        stop_checker: Optional[Callable[[], bool]] = None,
        max_steps: int = 50,
    ) -> Dict[str, Any]:
        """
        Run browser agent with checkpoint pattern.
        
        Browser agent has two tools:
        - mark_for_extraction: Save current page for seed extraction
        - checkpoint: Report status and get instructions from research agent
        
        Args:
            url: Starting URL
            goal: Initial goal for browser agent
            llm: LLM for browser agent
            checkpoint_callback: async fn(current_url, status) -> instructions
            extraction_dir: Where to save marked pages
            extraction_queue: Queue to put marked pages
            stop_checker: Optional fn() -> bool to check for pause
            max_steps: Max agent steps
            
        Returns:
            Dict with success, pages_marked, error
        """
        from browser_use import Agent, Tools
        
        extraction_dir.mkdir(parents=True, exist_ok=True)
        
        pages_marked = 0
        should_stop = False
        
        # Create tools
        tools = Tools()
        
        @tools.action(description="""Mark the current page for seed extraction.
Call this when you find a page with content that could become dataset rows.
The page will be saved and processed by the extraction system.
After marking, continue exploring or checkpoint for guidance.""")
        async def mark_for_extraction(description: str, browser_session) -> str:
            nonlocal pages_marked
            
            try:
                current_url = await browser_session.page.evaluate('() => window.location.href')
                html = await browser_session.page.evaluate('() => document.body.innerHTML')
                markdown = md(html, heading_style='ATX', strip=['script', 'style'])
                
                # Save to file
                filename = url_to_filename(current_url)
                filepath = extraction_dir / filename
                
                content = f"""---
source_url: {current_url}
extracted_at: {datetime.now(timezone.utc).isoformat()}
description: {description}
---

{markdown}
"""
                filepath.write_text(content, encoding='utf-8')
                
                # Queue for extraction
                await extraction_queue.put({
                    "file_path": str(filepath),
                    "source_url": current_url,
                    "description": description,
                })
                
                pages_marked += 1
                logger.info(f"[BrowserAgent] Marked page {pages_marked}: {filename}")
                
                return f"✓ Page marked for extraction ({pages_marked} total). Continue exploring or checkpoint."
                
            except Exception as e:
                logger.error(f"[BrowserAgent] mark_for_extraction failed: {e}")
                return f"Error marking page: {e}"
        
        @tools.action(description="""Checkpoint - report status and get instructions.
Call this when you:
- Reach a decision point
- Complete a section
- Need guidance on what to do next
- Are unsure whether to continue

Describe what you've found and any questions.""")
        async def checkpoint(status: str, browser_session) -> str:
            nonlocal should_stop
            
            if stop_checker and stop_checker():
                should_stop = True
                return "STOP - Task is being paused. Call done() to end."
            
            try:
                current_url = await browser_session.page.evaluate('() => window.location.href')
                instructions = await checkpoint_callback(current_url, status)
                
                logger.info(f"[BrowserAgent] Checkpoint at {current_url[:50]}...")
                logger.info(f"[BrowserAgent] Status: {status[:100]}...")
                logger.info(f"[BrowserAgent] Instructions: {instructions[:100]}...")
                
                # Check if research agent wants to end session
                instructions_lower = instructions.lower()
                if any(phrase in instructions_lower for phrase in ["end session", "stop", "done with this", "that's enough"]):
                    should_stop = True
                    return f"{instructions}\n\nCall done() to end the browser session."
                
                return instructions
                
            except Exception as e:
                logger.error(f"[BrowserAgent] checkpoint failed: {e}")
                return f"Checkpoint failed: {e}. Continue with your judgment or call done()."
        
        # Acquire browser and run agent
        await self._semaphore.acquire()
        browser = await self._available.get()
        browser_id = self._browser_ids.get(id(browser), "?")
        
        logger.info(f"[Browser-{browser_id}] Starting agent for: {url}")
        
        try:
            # Build task prompt
            task = f"""Navigate to {url} and accomplish this goal:

{goal}

## Your Tools
- mark_for_extraction(description): Save current page for seed extraction. Call when you find useful content.
- checkpoint(status): Report progress and get instructions. Call at decision points or when unsure.
- done(): End the session.

## Important
- Checkpoint frequently so the research coordinator can guide you
- Mark pages with content that could become dataset rows
- Don't try to extract data yourself - just mark pages
- If told to end session, call done()
"""
            
            agent = Agent(
                task=task,
                browser=browser,
                llm=llm,
                tools=tools,
                max_steps=max_steps,
            )
            
            history = await agent.run()
            
            success = history.is_successful() if hasattr(history, 'is_successful') else not should_stop
            error = history.errors()[-1] if hasattr(history, 'errors') and history.errors() else None
            
            return {
                "success": success,
                "pages_marked": pages_marked,
                "error": error,
            }
            
        except Exception as e:
            logger.error(f"[Browser-{browser_id}] Agent failed: {e}")
            return {
                "success": False,
                "pages_marked": pages_marked,
                "error": str(e),
            }
            
        finally:
            # Reset browser before returning to pool
            try:
                await browser.stop()
                await browser.start()
            except Exception as e:
                logger.warning(f"[Browser-{browser_id}] Reset failed: {e}")
            
            await self._available.put(browser)
            self._semaphore.release()
            logger.info(f"[Browser-{browser_id}] Released")